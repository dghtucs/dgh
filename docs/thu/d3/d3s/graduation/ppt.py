from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    
    # Define colors
    BG_COLOR = RGBColor(15, 23, 42)      # Dark Blue/Slate
    TEXT_MAIN = RGBColor(248, 250, 252)  # White/Slate
    ACCENT_CYAN = RGBColor(56, 189, 248) # Cyan
    ACCENT_INDIGO = RGBColor(129, 140, 248)# Indigo
    
    # Helper to set background (approximated by a rectangle)
    def set_background(slide):
        left = top = 0
        width = prs.slide_width
        height = prs.slide_height
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = BG_COLOR
        shape.line.fill.background() # No outline
        # Send to back hack: usually shapes added first are at back, but to be sure we just add content after.

    # Helper to add title
    def add_title(slide, text):
        title_shape = slide.shapes.title
        title_shape.text = text
        title_shape.text_frame.paragraphs[0].font.color.rgb = ACCENT_CYAN
        title_shape.text_frame.paragraphs[0].font.bold = True
        title_shape.text_frame.paragraphs[0].font.name = 'Arial'

    # Helper to add body text
    def add_text(slide, text, left, top, width, height, font_size=18, bold=False, color=TEXT_MAIN):
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = 'Arial'
        return tf

    # --- Slide 1: Title ---
    slide_layout = prs.slide_layouts[0] # Title Slide
    slide = prs.slides.add_slide(slide_layout)
    set_background(slide)
    
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "NetConfBench"
    title.text_frame.paragraphs[0].font.color.rgb = ACCENT_CYAN
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.size = Pt(60)
    
    subtitle.text = "A Holistic Framework to Evaluate LLM Agents\nfor Network Configuration"
    subtitle.text_frame.paragraphs[0].font.color.rgb = TEXT_MAIN
    subtitle.text_frame.paragraphs[0].font.size = Pt(32)

    # --- Slide 2: Background & Motivation ---
    slide_layout = prs.slide_layouts[1] # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    set_background(slide)
    add_title(slide, "Background & Motivation")
    
    # Left Column: Challenge
    add_text(slide, "The Challenge:", Inches(0.5), Inches(1.5), Inches(4), Inches(0.5), font_size=24, bold=True, color=ACCENT_INDIGO)
    content_l = add_text(slide, "", Inches(0.5), Inches(2.2), Inches(4), Inches(4))
    p = content_l.add_paragraph()
    p.text = "• Network configuration ensures stability but relies on human expertise."
    p.font.size = Pt(18); p.font.color.rgb = TEXT_MAIN
    p = content_l.add_paragraph()
    p.text = "• Manual config is error-prone and slow."
    p.font.size = Pt(18); p.font.color.rgb = TEXT_MAIN

    # Right Column: Opportunity
    add_text(slide, "The Opportunity:", Inches(5), Inches(1.5), Inches(4), Inches(0.5), font_size=24, bold=True, color=ACCENT_INDIGO)
    content_r = add_text(slide, "", Inches(5), Inches(2.2), Inches(4), Inches(4))
    p = content_r.add_paragraph()
    p.text = "• LLMs show promise in semantic understanding."
    p.font.size = Pt(18); p.font.color.rgb = TEXT_MAIN
    p = content_r.add_paragraph()
    p.text = "• Research Question: Can LLMs function as autonomous agents?"
    p.font.size = Pt(20); p.font.bold = True; p.font.color.rgb = ACCENT_CYAN

    # --- Slide 3: Limitations of Existing Benchmarks ---
    slide = prs.slides.add_slide(prs.slide_layouts[5]) # Blank
    set_background(slide)
    add_text(slide, "Limitations of Existing Benchmarks", Inches(0.5), Inches(0.5), Inches(9), Inches(1), font_size=40, bold=True, color=ACCENT_CYAN)

    # Table
    rows = 4
    cols = 3
    left = Inches(1.0)
    top = Inches(2.0)
    width = Inches(8.0)
    height = Inches(3.0)
    
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    
    # Headers
    headers = ["Feature", "Existing Benchmarks", "NetConfBench (Ours)"]
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(0,0,0) # Black text for table clarity if default is light
        
    data = [
        ["Environment", "Static / Text-based", "Dynamic / Emulated (GNS3)"],
        ["Task Complexity", "Simple / Synthetic", "Realistic / Multi-step"],
        ["Interaction", "One-way Generation", "Closed-loop (Perception-Action)"]
    ]
    
    for r_idx, row_data in enumerate(data):
        for c_idx, val in enumerate(row_data):
            cell = table.cell(r_idx+1, c_idx)
            cell.text = val
            cell.text_frame.paragraphs[0].font.size = Pt(16)
            if c_idx == 2:
                cell.text_frame.paragraphs[0].font.bold = True
                cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 100, 0) # Dark Green

    # --- Slide 4: Framework Overview ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_background(slide)
    add_title(slide, "NetConfBench Framework")
    
    # 3 Columns
    col_width = Inches(3)
    y_pos = Inches(2.5)
    
    # 1
    add_text(slide, "1. Task Preparation", Inches(0.5), y_pos, col_width, Inches(0.5), font_size=20, bold=True, color=ACCENT_INDIGO)
    add_text(slide, "Selection from dataset. Auto-initialization of topology and base configs.", Inches(0.5), y_pos + Inches(0.7), col_width, Inches(2), font_size=16)

    # 2
    add_text(slide, "2. Task Execution", Inches(3.75), y_pos, col_width, Inches(0.5), font_size=20, bold=True, color=ACCENT_CYAN)
    add_text(slide, "Agent interacts via APIs: Perception, Decision-making, and Action.", Inches(3.75), y_pos + Inches(0.7), col_width, Inches(2), font_size=16)

    # 3
    add_text(slide, "3. Task Evaluation", Inches(7), y_pos, col_width, Inches(0.5), font_size=20, bold=True, color=ACCENT_INDIGO)
    add_text(slide, "Assessment of reasoning traces, config correctness, and validation.", Inches(7), y_pos + Inches(0.7), col_width, Inches(2), font_size=16)

    # --- Slide 5: Dataset ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_background(slide)
    add_title(slide, "Benchmark Dataset Construction")
    
    content = add_text(slide, "", Inches(0.5), Inches(1.5), Inches(4), Inches(4))
    content.add_paragraph().text = "Dataset Composition:"
    p = content.add_paragraph(); p.text = "• 40 Tasks (Diverse Scenarios)"; p.font.size = Pt(18); p.font.color.rgb = TEXT_MAIN
    p = content.add_paragraph(); p.text = "• Routing (OSPF, BGP)"; p.font.size = Pt(18); p.font.color.rgb = TEXT_MAIN
    p = content.add_paragraph(); p.text = "• QoS (Traffic Shaping)"; p.font.size = Pt(18); p.font.color.rgb = TEXT_MAIN
    p = content.add_paragraph(); p.text = "• Security (ACLs)"; p.font.size = Pt(18); p.font.color.rgb = TEXT_MAIN

    content_r = add_text(slide, "Pipeline:", Inches(5), Inches(1.5), Inches(4.5), Inches(4))
    p = content_r.add_paragraph(); p.text = "1. Data Collection (Community labs)"; p.font.size = Pt(18); p.font.color.rgb = TEXT_MAIN
    p = content_r.add_paragraph(); p.text = "2. Standardization (Unified JSON)"; p.font.size = Pt(18); p.font.color.rgb = TEXT_MAIN
    p = content_r.add_paragraph(); p.text = "3. Reasoning Gen (GPT-o3 Ground Truth)"; p.font.size = Pt(18); p.font.color.rgb = TEXT_MAIN
    p = content_r.add_paragraph(); p.text = "4. Expert Verification"; p.font.size = Pt(18); p.font.bold=True; p.font.color.rgb = ACCENT_CYAN

    # --- Slide 6: Environment ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_background(slide)
    add_title(slide, "Interactive Environment (GNS3)")
    
    add_text(slide, "Standardized APIs for Agent Interaction:", Inches(0.5), Inches(1.5), Inches(9), Inches(0.5), font_size=20)
    
    box_w = Inches(4)
    box_h = Inches(1.5)
    
    # Box 1
    add_text(slide, "get-topology", Inches(0.5), Inches(2.5), box_w, Inches(0.5), font_size=20, bold=True, color=ACCENT_CYAN)
    add_text(slide, "Inspect nodes, links, and interface states.", Inches(0.5), Inches(3.0), box_w, Inches(1), font_size=16)

    # Box 2
    add_text(slide, "get-running-config", Inches(5), Inches(2.5), box_w, Inches(0.5), font_size=20, bold=True, color=ACCENT_CYAN)
    add_text(slide, "Read current device configurations.", Inches(5), Inches(3.0), box_w, Inches(1), font_size=16)

    # Box 3
    add_text(slide, "update-config", Inches(0.5), Inches(4.5), box_w, Inches(0.5), font_size=20, bold=True, color=ACCENT_CYAN)
    add_text(slide, "Apply configuration commands to devices.", Inches(0.5), Inches(5.0), box_w, Inches(1), font_size=16)

    # Box 4
    add_text(slide, "execute-validation", Inches(5), Inches(4.5), box_w, Inches(0.5), font_size=20, bold=True, color=ACCENT_CYAN)
    add_text(slide, "Run verification checks (ping, show commands).", Inches(5), Inches(5.0), box_w, Inches(1), font_size=16)

    # --- Slide 7: Metrics ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_background(slide)
    add_title(slide, "Multi-dimensional Evaluation")
    
    col_width = Inches(2.8)
    spacing = Inches(0.2)
    start_x = Inches(0.5)
    
    # 1
    add_text(slide, "Reasoning Score", start_x, Inches(2.5), col_width, Inches(0.5), font_size=22, bold=True, color=ACCENT_INDIGO)
    add_text(slide, "Cosine Similarity.\nCompares Agent thought process vs Ground Truth.", start_x, Inches(3.2), col_width, Inches(2), font_size=16)
    
    # 2
    add_text(slide, "Command Score", start_x + col_width + spacing, Inches(2.5), col_width, Inches(0.5), font_size=22, bold=True, color=ACCENT_INDIGO)
    add_text(slide, "Precision & Recall.\nChecks if config lines match solution set.", start_x + col_width + spacing, Inches(3.2), col_width, Inches(2), font_size=16)

    # 3
    add_text(slide, "Testcase Score", start_x + (col_width + spacing)*2, Inches(2.5), col_width, Inches(0.5), font_size=22, bold=True, color=ACCENT_CYAN)
    add_text(slide, "Functional Pass Rate.\nExecutes validation scripts. The most critical metric.", start_x + (col_width + spacing)*2, Inches(3.2), col_width, Inches(2), font_size=16)

    # --- Slide 8: Setup ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_background(slide)
    add_title(slide, "Experimental Setup")
    
    content = add_text(slide, "", Inches(1), Inches(2), Inches(8), Inches(4))
    p = content.add_paragraph(); p.text = "Tested Models (6 Representative LLMs):"; p.font.bold=True; p.font.size=Pt(24); p.font.color.rgb=ACCENT_CYAN
    p = content.add_paragraph(); p.text = "• Reasoning Models: GPT-o3, DeepSeek-r1, QwQ-plus"; p.font.size=Pt(20); p.font.color.rgb=TEXT_MAIN
    p = content.add_paragraph(); p.text = "• Non-Reasoning Models: GPT-4o, DeepSeek-v3, Qwen-max"; p.font.size=Pt(20); p.font.color.rgb=TEXT_MAIN
    
    content.add_paragraph().text = "" # spacer
    
    p = content.add_paragraph(); p.text = "Interaction Modes:"; p.font.bold=True; p.font.size=Pt(24); p.font.color.rgb=ACCENT_CYAN
    p = content.add_paragraph(); p.text = "1. One-shot Prediction (All models)"; p.font.size=Pt(20); p.font.color.rgb=TEXT_MAIN
    p = content.add_paragraph(); p.text = "2. Multi-turn Interaction (ReAct for non-reasoning)"; p.font.size=Pt(20); p.font.color.rgb=TEXT_MAIN

    # --- Slide 9: RQ1 ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_background(slide)
    add_title(slide, "RQ1: Reasoning Capability Analysis")
    
    add_text(slide, "Conclusion: Advanced reasoning capability is the primary differentiator.", Inches(1), Inches(1.5), Inches(8), Inches(1), font_size=20, color=TEXT_MAIN)
    
    # Textual Chart
    chart_txt = add_text(slide, "", Inches(1), Inches(2.5), Inches(8), Inches(3))
    data_pts = [
        ("GPT-o3", "90.2%", ACCENT_CYAN),
        ("DeepSeek-r1", "82.4%", ACCENT_CYAN),
        ("GPT-4o", "71.4%", TEXT_MAIN),
        ("Qwen-max", "58.0%", TEXT_MAIN),
        ("QwQ-plus", "48.7%", RGBColor(239, 68, 68))
    ]
    for name, val, color in data_pts:
        p = chart_txt.add_paragraph()
        p.text = f"{name}: {val}"
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = color
    
    add_text(slide, "* Metric: Average Testcase Pass Rate", Inches(1), Inches(6), Inches(5), Inches(0.5), font_size=14, color=RGBColor(148, 163, 184))

    # --- Slide 10: RQ2 ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_background(slide)
    add_title(slide, "RQ2: Impact of Task Complexity")
    
    content_l = add_text(slide, "", Inches(0.5), Inches(2), Inches(4.5), Inches(4))
    p = content_l.add_paragraph(); p.text = "Global Planning Deficit"; p.font.bold=True; p.font.size=Pt(24); p.font.color.rgb=ACCENT_CYAN
    p = content_l.add_paragraph(); p.text = "Models perform well on isolated subtasks but struggle when integrated into complex tasks."; p.font.size=Pt(20); p.font.color.rgb=TEXT_MAIN
    
    content_r = add_text(slide, "", Inches(5.5), Inches(2), Inches(4), Inches(4))
    p = content_r.add_paragraph(); p.text = "Pass Rate Comparison:"; p.font.bold=True; p.font.size=Pt(20); p.font.color.rgb=TEXT_MAIN
    p = content_r.add_paragraph(); p.text = "Subtasks: ~85%"; p.font.bold=True; p.font.size=Pt(36); p.font.color.rgb=RGBColor(52, 211, 153)
    p = content_r.add_paragraph(); p.text = "Full Task: ~40%"; p.font.bold=True; p.font.size=Pt(36); p.font.color.rgb=RGBColor(251, 191, 36)

    # --- Slide 11: RQ3 ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_background(slide)
    add_title(slide, "RQ3: Effect of Multi-turn Interaction")
    
    content_l = add_text(slide, "", Inches(0.5), Inches(2), Inches(4.5), Inches(4))
    p = content_l.add_paragraph(); p.text = "Strengths"; p.font.bold=True; p.font.size=Pt(24); p.font.color.rgb=RGBColor(52, 211, 153)
    p = content_l.add_paragraph(); p.text = "• Corrects Syntax Errors"; p.font.size=Pt(20); p.font.color.rgb=TEXT_MAIN
    p = content_l.add_paragraph(); p.text = "• Adjusts Interface Modes"; p.font.size=Pt(20); p.font.color.rgb=TEXT_MAIN

    content_r = add_text(slide, "", Inches(5.5), Inches(2), Inches(4.5), Inches(4))
    p = content_r.add_paragraph(); p.text = "Weaknesses"; p.font.bold=True; p.font.size=Pt(24); p.font.color.rgb=RGBColor(251, 191, 36)
    p = content_r.add_paragraph(); p.text = "• Cannot solve high-level planning issues"; p.font.size=Pt(20); p.font.color.rgb=TEXT_MAIN
    p = content_r.add_paragraph(); p.text = "• Self-Verification Failure (Hallucinating success)"; p.font.size=Pt(20); p.font.color.rgb=TEXT_MAIN

    # --- Slide 12: Conclusion ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_background(slide)
    add_title(slide, "Conclusion & Future Work")
    
    content = add_text(slide, "", Inches(1), Inches(2), Inches(8), Inches(4))
    p = content.add_paragraph(); p.text = "Summary:"; p.font.bold=True; p.font.size=Pt(24); p.font.color.rgb=ACCENT_CYAN
    p = content.add_paragraph(); p.text = "NetConfBench is the first holistic framework. While Reasoning models show promise, robustness in production is still lacking."; p.font.size=Pt(20); p.font.color.rgb=TEXT_MAIN
    
    content.add_paragraph().text = ""
    
    p = content.add_paragraph(); p.text = "Future Directions:"; p.font.bold=True; p.font.size=Pt(24); p.font.color.rgb=ACCENT_INDIGO
    p = content.add_paragraph(); p.text = "• Reinforcement Learning (RL) training ground"; p.font.size=Pt(20); p.font.color.rgb=TEXT_MAIN
    p = content.add_paragraph(); p.text = "• Fault Diagnosis and Troubleshooting"; p.font.size=Pt(20); p.font.color.rgb=TEXT_MAIN

    # Save
    prs.save('NetConfBench_Presentation.pptx')
    print("Presentation saved as NetConfBench_Presentation.pptx")

if __name__ == "__main__":
    create_presentation()